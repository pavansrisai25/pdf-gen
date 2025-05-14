import streamlit as st
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from io import BytesIO
from PIL import Image
import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import zipfile
import base64
st.title("PDF-GENERATOR")
# Predefined scrolling text
ticker_text = (
    "💡  Transforming PDFs with ease—convert, merge, split, and more! | "
    "🎉 100% Free & Fast—No subscriptions, No limits! | "
    "📂 Upload, Process, Download—All in just a few clicks! | "
    "💚 Trusted by thousands—Join the community today! | "
    "🔗 Streamlit-powered for smooth and secure file handling! | "
    "📢 Stay tuned for exciting updates & new features! | "
)
# HTML + JavaScript for the scrolling ticker
ticker_html = f"""
<iframe width="100%" height="50" frameborder="0" scrolling="no" style="border:none; overflow:hidden;"
srcdoc='
    <style>
        @keyframes ticker {{
            from {{ transform: translateX(100%); }}
            to {{ transform: translateX(-100%); }}
        }}
        .ticker-container {{
            width: 100%;
            overflow: hidden;
            background: #008000;
            color: white;
            font-size: 20px;
            font-weight: bold;
            white-space: nowrap;
            position: fixed;
            bottom: 0;
            left: 0;
            padding: 10px 0;
            z-index: 1000;
        }}
        .ticker-text {{
            display: inline-block;
            padding-left: 100%;
            animation: ticker 50s linear infinite;
        }}
    </style>
    <div class="ticker-container">
        <div class="ticker-text">{ticker_text}</div>
    </div>
'>
</iframe>
"""

# Display the moving ticker
st.markdown(ticker_html, unsafe_allow_html=True)
# Display the logo with auto-width adjustment
st.image("logo1.png", width=150)

# Function to set background image from URL
def set_bg_from_url():
    st.markdown(
        """
        <style>
        .stApp {
            /* Replace this URL with your raw GitHub link */
            background-image: url("https://raw.githubusercontent.com/Siva369IT/Pdf-well/main/bgimage.png");
            background-size: cover;
            background-position: center;
            background-repeat: no-repeat;
            background-attachment: fixed;
        }
        /* Mobile screens: auto-fit the image height */
        @media (max-width: 768px) {
            .stApp {
                background-size: auto 100%;
            }
        }
        </style>
        """,
        unsafe_allow_html=True
    )

# Apply the background styling
set_bg_from_url()
operation = st.selectbox(
    "Choose an Operation",
    (
        "Click here to choose 👇",
        "Generate Empty PDF",
        "Convert Any File to PDF",
        "Extract Pages from PDF",
        "Merge PDFs",
        "Split PDF",
        "Compress PDF",
        "Insert Page Numbers",
        "Images to PDF",
        "Remove uploaded files"
    )
)

# Dynamic Upload Instructions
if operation == "Convert Any File to PDF":
    st.info("Upload .txt, .doc, .docx, .ppt, .pptx, or image files (jpg, png, jpeg)")
elif operation == "Extract Pages from PDF" or operation in ["Merge PDFs", "Split PDF", "Compress PDF", "Insert Page Numbers"]:
    st.info("Upload PDF files")
elif operation == "Images to PDF":
    st.info("Upload multiple image files (jpg, png, jpeg)")
elif operation == "Generate Empty PDF":
    st.info("No file needed. Just enter number of pages.")
else:
    st.info("Welcome to pdf-world 💚")

# Uploads
if operation in ["Convert Any File to PDF", "Images to PDF"]:
    uploaded_files = st.file_uploader("Upload Files", type=['txt', 'doc', 'docx', 'ppt', 'pptx', 'png', 'jpg', 'jpeg'], accept_multiple_files=True)
elif operation in ["Extract Pages from PDF", "Compress PDF", "Split PDF", "Insert Page Numbers"]:
    uploaded_file = st.file_uploader("Upload PDF file", type='pdf')
elif operation == "Merge PDFs":
    uploaded_files = st.file_uploader("Upload exactly 2 PDF files", type='pdf', accept_multiple_files=True)
else:
    uploaded_file = None

def download_button(file, file_name):
    st.download_button(label="Download", data=file, file_name=file_name, mime="application/pdf")


# 1. Generate Empty PDF
if operation == "Generate Empty PDF":
    st.subheader("Generate Empty PDF")
    num_pages = st.number_input("Enter the number of pages to generate (max 3690):", min_value=1, max_value=3690, step=1)
    if st.button("Generate"):
        output_pdf = BytesIO()
        pdf_writer = PdfWriter()
        for _ in range(num_pages):
            pdf_writer.add_blank_page(width=612, height=792)
        pdf_writer.write(output_pdf)
        output_pdf.seek(0)
        st.success(f"Generated Empty PDF with {num_pages} pages.")
        custom_name = st.text_input("Enter custom file name (without extension):", value="empty_pdf")
        st.download_button(
            label="Download Empty PDF",
            data=output_pdf,
            file_name=f"{custom_name}.pdf",
            mime="application/pdf"
        )
# 2. Convert Any File to PDF (direct download)
elif operation == "Convert Any File to PDF" and uploaded_files:
    for uploaded in uploaded_files:
        filename = uploaded.name
        file_bytes = uploaded.read()
        file_ext = filename.split('.')[-1].lower()
        output = BytesIO()

        try:
            if file_ext == "txt":
                c = canvas.Canvas(output, pagesize=letter)
                text = file_bytes.decode("utf-8")
                text_lines = text.split('\n')
                y = 750
                for line in text_lines:
                    c.drawString(72, y, line)
                    y -= 15
                    if y < 72:
                        c.showPage()
                        y = 750
                c.save()
            elif file_ext in ["doc", "docx"]:
                doc = Document(BytesIO(file_bytes))
                c = canvas.Canvas(output, pagesize=letter)
                y = 750
                for para in doc.paragraphs:
                    c.drawString(72, y, para.text)
                    y -= 15
                    if y < 72:
                        c.showPage()
                        y = 750
                c.save()
            elif file_ext in ["ppt", "pptx"]:
                presentation = Presentation(BytesIO(file_bytes))
                c = canvas.Canvas(output, pagesize=letter)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            c.drawString(72, 750, shape.text)
                            c.showPage()
                c.save()
            elif file_ext in ["png", "jpg", "jpeg"]:
                img = Image.open(BytesIO(file_bytes)).convert("RGB")
                pdf_buffer = BytesIO()
                img.save(pdf_buffer, format="PDF")
                pdf_buffer.seek(0)
                output = pdf_buffer
            else:
                st.warning(f"Incorrect format ❗: {file_ext}")
                continue

            output.seek(0)
            download_button(output, f"{filename.split('.')[0]}.pdf")
        except Exception as e:
            st.error(f"Error converting {filename}: {e}")

# 3. Extract Pages
elif operation == "Extract Pages from PDF" and uploaded_file:
    page_input = st.text_input("Enter page numbers or ranges (e.g., 1,3,5-8):")
    if st.button("Extract Pages"):
        try:
            reader = PdfReader(uploaded_file)
            writer = PdfWriter()
            ranges = page_input.replace(' ', '').split(',')
            pages_to_extract = []
            for r in ranges:
                if '-' in r:
                    start, end = map(int, r.split('-'))
                    pages_to_extract.extend(range(start-1, end))
                else:
                    pages_to_extract.append(int(r)-1)
            invalid_pages = [p+1 for p in pages_to_extract if p >= len(reader.pages)]
            if invalid_pages:
                st.error(f"Pages not found: {invalid_pages}")
            else:
                for p in pages_to_extract:
                    writer.add_page(reader.pages[p])
                output = BytesIO()
                writer.write(output)
                output.seek(0)
                download_button(output, "extracted_pages.pdf")
        except Exception as e:
            st.error(f"Extraction error: {e}")

# 4. Merge PDFs
elif operation == "Merge PDFs" and uploaded_files:
    if len(uploaded_files) == 2:
        merger = PdfWriter()
        for pdf in uploaded_files:
            merger.append(PdfReader(pdf))
        output = BytesIO()
        merger.write(output)
        output.seek(0)
        download_button(output, "merged.pdf")
    else:
        st.warning("Please upload exactly 2 PDF files.")

# 5. Split PDF
elif operation == "Split PDF" and uploaded_file:
    reader = PdfReader(uploaded_file)
    choice = st.radio("Choose split method:", ("Custom Split", "Split Each Page"))
    if choice == "Custom Split":
        split_at = st.number_input("Split after how many pages?", min_value=1, max_value=len(reader.pages)-1)
        if st.button("Split PDF"):
            part1_writer = PdfWriter()
            part2_writer = PdfWriter()
            for i, page in enumerate(reader.pages):
                if i < split_at:
                    part1_writer.add_page(page)
                else:
                    part2_writer.add_page(page)
            out1, out2 = BytesIO(), BytesIO()
            part1_writer.write(out1)
            part2_writer.write(out2)
            out1.seek(0)
            out2.seek(0)
            st.download_button("Download Part 1", data=out1, file_name="part1.pdf")
            st.download_button("Download Part 2", data=out2, file_name="part2.pdf")
    else:
        if st.button("Split into single-page PDFs"):
            
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w') as zipf:
                for i, page in enumerate(reader.pages):
                    writer = PdfWriter()
                    writer.add_page(page)
                    part_buffer = BytesIO()
                    writer.write(part_buffer)
                    part_buffer.seek(0)
                    zipf.writestr(f"page_{i+1}.pdf", part_buffer.read())
            zip_buffer.seek(0)
            st.download_button("Download ZIP", data=zip_buffer, file_name="split_pages.zip")

# 6. Compress PDF
elif operation == "Compress PDF" and uploaded_file:
    compress_ratio = st.slider("Compression level", 1, 100, 50)
    if st.button("Compress"):
        try:
            input_pdf = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            output_pdf = fitz.open()
            for page in input_pdf:
                pix = page.get_pixmap(matrix=fitz.Matrix(1, 1))
                img_bytes = pix.tobytes("jpeg")
                new_page = output_pdf.new_page(width=page.rect.width, height=page.rect.height)
                new_page.insert_image(new_page.rect, stream=img_bytes)
            output = BytesIO()
            output_pdf.save(output)
            output.seek(0)
            download_button(output, "compressed.pdf")
            st.info("Note: Compression depends on PDF content. Some files may not reduce in size.")
        except Exception as e:
            st.error(f"Compression error: {e}")

# 7. Insert Page Numbers
elif operation == "Insert Page Numbers" and uploaded_file:
    if st.button("Insert Page Numbers"):
        try:
            reader = PdfReader(uploaded_file)
            writer = PdfWriter()
            for i, page in enumerate(reader.pages):
                packet = BytesIO()
                c = canvas.Canvas(packet, pagesize=letter)
                c.drawString(300, 15, str(i+1))
                c.save()
                packet.seek(0)
                new_pdf = PdfReader(packet)
                page.merge_page(new_pdf.pages[0])
                writer.add_page(page)
            output = BytesIO()
            writer.write(output)
            output.seek(0)
            download_button(output, "with_page_numbers.pdf")
        except Exception as e:
            st.error(f"Insertion error: {e}")

# 8. Images to PDF
elif operation == "Images to PDF" and uploaded_files:
    if st.button("Convert Images to Single PDF"):
        try:
            images = [Image.open(f).convert("RGB") for f in uploaded_files]
            output = BytesIO()
            images[0].save(output, save_all=True, append_images=images[1:], format="PDF")
            output.seek(0)
            download_button(output, "images_combined.pdf")
        except Exception as e:
            st.error(f"Image conversion error: {e}")

# Remove all
elif operation == "Remove Uploaded Files":
    st.stop()

# Footer at the bottom

st.markdown(
"""<div style="position: fixed; bottom: 45px; width: 100%; text-align: center; font-size: small; color: #2F4F4F;">
© 2025 Pavan SriSai Mondem | Siva Satyamsetti | Uma Satya Mounika Sapireddy | Bhuvaneswari Devi Seru | Chandu Meela | Trainees from techwing 🧡 | Made with 💚 by gen ai and aws trainees

   </div>""",  
    unsafe_allow_html=True  
)
